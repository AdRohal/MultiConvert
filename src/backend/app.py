from flask import Flask, request, send_file, jsonify
from werkzeug.utils import secure_filename
from io import BytesIO
import os
import logging
import fitz  # PyMuPDF
from docx import Document
from docx.shared import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
import tempfile  # Import tempfile module
from flask_cors import CORS  # Import CORS
from docx2pdf import convert  # Import docx2pdf
from pypandoc import convert_file  # Import pypandoc
from pptx import Presentation  # Import python-pptx
from pptx.util import Inches, Pt  # Import pptx.util for setting slide dimensions
import comtypes.client  # Import comtypes for converting pptx to pdf
import pythoncom  # Import pythoncom for COM initialization
import zipfile  # Import zipfile for creating ZIP files
from pydub import AudioSegment  # Import pydub for audio conversion
from PIL import Image  # Import PIL for image conversion

app = Flask(__name__)
CORS(app)  # Enable CORS

# Configure logging
logging.basicConfig(level=logging.DEBUG)

# Directory to save converted files
CONVERTED_FOLDER = os.path.join(os.path.dirname(__file__), 'converted_files')
os.makedirs(CONVERTED_FOLDER, exist_ok=True)

@app.route('/convert', methods=['POST'])
def convert_file():
    try:
        if 'file' not in request.files:
            app.logger.error("No file part in the request")
            return jsonify({"error": "No file part"}), 400

        file = request.files['file']
        format = request.form.get('format', 'docx').lower()

        if file.filename == '':
            app.logger.error("No selected file")
            return jsonify({"error": "No selected file"}), 400

        filename = secure_filename(file.filename)
        file_ext = filename.rsplit('.', 1)[1].lower()

        app.logger.debug(f"Received file: {filename}, format: {format}")

        if file_ext in ['png', 'jpg', 'jpeg'] and format in ['png', 'jpg', 'jpeg', 'ico']:
            # Convert image to the desired format
            image = Image.open(file)
            converted_filename = f"{filename.rsplit('.', 1)[0]}.{format}"
            converted_file_path = os.path.join(CONVERTED_FOLDER, converted_filename)
            if format in ['jpg', 'jpeg']:
                image = image.convert("RGB")  # Ensure the image is in RGB mode for JPG conversion
            image.save(converted_file_path, format=format.upper() if format != 'ico' else 'ICO')

            app.logger.info(f"Successfully converted {filename} to {format.upper()}")
            return jsonify({"filename": converted_filename})

        elif file_ext == 'png' and format == 'pdf':
            # Convert PNG to PDF
            image = Image.open(file)
            pdf_filename = f"{filename.rsplit('.', 1)[0]}.pdf"
            pdf_file_path = os.path.join(CONVERTED_FOLDER, pdf_filename)

            image = image.convert("RGB")  # Ensure the image is in RGB mode for PDF conversion
            image.save(pdf_file_path, "PDF", resolution=100.0)

            app.logger.info(f"Successfully converted {filename} to PDF")
            return jsonify({"filename": pdf_filename})

        elif file_ext == 'png' and format == 'docx':
            # Convert PNG to DOCX
            image = Image.open(file)
            docx_document = Document()

            # Set page size to match image size
            width, height = image.size
            section = docx_document.sections[-1]
            section.page_width = Pt(width)
            section.page_height = Pt(height)
            section.left_margin = Inches(0)
            section.right_margin = Inches(0)
            section.top_margin = Inches(0)
            section.bottom_margin = Inches(0)

            # Add image to DOCX
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_img:
                image.save(temp_img, format="PNG")
                temp_img_path = temp_img.name

            docx_document.add_picture(temp_img_path, width=Pt(width), height=Pt(height))
            os.remove(temp_img_path)

            # Save the DOCX file to the designated directory
            docx_filename = f"{filename.rsplit('.', 1)[0]}.docx"
            docx_file_path = os.path.join(CONVERTED_FOLDER, docx_filename)
            docx_document.save(docx_file_path)

            app.logger.info(f"Successfully converted {filename} to DOCX")
            return jsonify({"filename": docx_filename})

        elif file_ext == 'pdf' and format == 'docx':
            # Convert PDF to DOCX by rendering each page as an image
            pdf_file = BytesIO(file.read())
            pdf_document = fitz.open(stream=pdf_file, filetype="pdf")
            docx_document = Document()

            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                pix = page.get_pixmap()

                # Save page as an image
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_img:
                    temp_img.write(pix.tobytes())
                    temp_img_path = temp_img.name

                # Set page size to match PDF page size
                pdf_width, pdf_height = page.rect.width, page.rect.height
                section = docx_document.sections[-1]
                section.page_width = Pt(pdf_width)
                section.page_height = Pt(pdf_height)
                section.left_margin = Inches(0)
                section.right_margin = Inches(0)
                section.top_margin = Inches(0)
                section.bottom_margin = Inches(0)

                # Add image to DOCX
                docx_document.add_picture(temp_img_path, width=section.page_width, height=section.page_height)
                os.remove(temp_img_path)

                # Add a page break after each page except the last one
                if page_num < len(pdf_document) - 1:
                    docx_document.add_section()

            # Save the DOCX file to the designated directory
            converted_filename = f"{filename.rsplit('.', 1)[0]}.docx"
            converted_file_path = os.path.join(CONVERTED_FOLDER, converted_filename)
            docx_document.save(converted_file_path)

            app.logger.info(f"Successfully converted {filename} to DOCX")
            return jsonify({"filename": converted_filename})
        
        elif file_ext == 'docx' and format == 'pdf':
            # Convert DOCX to PDF using docx2pdf
            docx_file = BytesIO(file.read())
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_docx:
                temp_docx.write(docx_file.read())
                temp_docx_path = temp_docx.name

            pdf_filename = f"{filename.rsplit('.', 1)[0]}.pdf"
            pdf_file_path = os.path.join(CONVERTED_FOLDER, pdf_filename)

            pythoncom.CoInitialize()  # Initialize COM library
            try:
                # Convert DOCX to PDF
                convert(temp_docx_path, pdf_file_path)
            finally:
                pythoncom.CoUninitialize()  # Uninitialize COM library

            app.logger.info(f"Successfully converted {filename} to PDF")
            return jsonify({"filename": pdf_filename})
        
        elif file_ext == 'docx' and format == 'pptx':
            # Convert DOCX to PPTX using python-pptx
            docx_file = BytesIO(file.read())
            docx_document = Document(docx_file)
            pptx_document = Presentation()

            for para in docx_document.paragraphs:
                slide_layout = pptx_document.slide_layouts[5]  # Use a blank slide layout
                slide = pptx_document.slides.add_slide(slide_layout)
                text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8.5), Inches(5.5))
                text_frame = text_box.text_frame
                p = text_frame.add_paragraph()
                p.text = para.text

            # Save the PPTX file to the designated directory
            pptx_filename = f"{filename.rsplit('.', 1)[0]}.pptx"
            pptx_file_path = os.path.join(CONVERTED_FOLDER, pptx_filename)
            pptx_document.save(pptx_file_path)

            app.logger.info(f"Successfully converted {filename} to PPTX")
            return jsonify({"filename": pptx_filename})
        
        elif file_ext == 'pdf' and format == 'doc':
            # Convert PDF to DOC using pypandoc
            pdf_file = BytesIO(file.read())
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
                temp_pdf.write(pdf_file.read())
                temp_pdf_path = temp_pdf.name

            doc_filename = f"{filename.rsplit('.', 1)[0]}.doc"
            doc_file_path = os.path.join(CONVERTED_FOLDER, doc_filename)

            # Convert PDF to DOC
            convert_file(temp_pdf_path, 'doc', outputfile=doc_file_path)

            app.logger.info(f"Successfully converted {filename} to DOC")
            return jsonify({"filename": doc_filename})
        
        elif file_ext == 'pdf' and format in ['png', 'jpg', 'jpeg']:
            # Convert PDF to PNG or JPG
            pdf_file = BytesIO(file.read())
            pdf_document = fitz.open(stream=pdf_file, filetype="pdf")
            image_paths = []

            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                pix = page.get_pixmap()

                # Save page as an image
                image_filename = f"{filename.rsplit('.', 1)[0]}_page_{page_num + 1}.{format}"
                image_path = os.path.join(CONVERTED_FOLDER, image_filename)
                pix.save(image_path)
                image_paths.append(image_path)

            if len(image_paths) > 1:
                # Create a ZIP file if there are multiple images
                zip_filename = f"{filename.rsplit('.', 1)[0]}_images.zip"
                zip_path = os.path.join(CONVERTED_FOLDER, zip_filename)
                with zipfile.ZipFile(zip_path, 'w') as zipf:
                    for image_path in image_paths:
                        zipf.write(image_path, os.path.basename(image_path))
                        os.remove(image_path)  # Remove the image file after adding to ZIP

                app.logger.info(f"Successfully converted {filename} to ZIP containing images")
                return jsonify({"filename": zip_filename})
            else:
                app.logger.info(f"Successfully converted {filename} to {format.upper()}")
                return jsonify({"filename": os.path.basename(image_paths[0])})
        
        elif file_ext == 'pptx' and format == 'pdf':
            # Convert PPTX to PDF using comtypes
            pptx_file = BytesIO(file.read())
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_pptx:
                temp_pptx.write(pptx_file.read())
                temp_pptx_path = temp_pptx.name

            pdf_filename = f"{filename.rsplit('.', 1)[0]}.pdf"
            pdf_file_path = os.path.join(CONVERTED_FOLDER, pdf_filename)

            pythoncom.CoInitialize()  # Initialize COM library
            try:
                powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                powerpoint.Visible = 1
                presentation = powerpoint.Presentations.Open(temp_pptx_path)
                presentation.SaveAs(pdf_file_path, 32)  # 32 represents the format for PDF
                presentation.Close()
                powerpoint.Quit()
            finally:
                pythoncom.CoUninitialize()  # Uninitialize COM library

            app.logger.info(f"Successfully converted {filename} to PDF")
            return jsonify({"filename": pdf_filename})
        
        elif file_ext == 'pdf' and format == 'pdf':
            # Re-save PDF
            pdf_file = BytesIO(file.read())
            pdf_filename = f"{filename.rsplit('.', 1)[0]}_converted.pdf"
            pdf_file_path = os.path.join(CONVERTED_FOLDER, pdf_filename)

            with open(pdf_file_path, 'wb') as f:
                f.write(pdf_file.read())

            app.logger.info(f"Successfully re-saved {filename} as PDF")
            return jsonify({"filename": pdf_filename})
        
        else:
            app.logger.error("Unsupported file format or conversion")
            return jsonify({"error": "Unsupported file format or conversion"}), 400
    except Exception as e:
        app.logger.error(f"Error during file conversion: {e}")
        return jsonify({"error": f"An error occurred during file conversion: {e}"}), 500

@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    file_path = os.path.join(CONVERTED_FOLDER, filename)
    app.logger.debug(f"Attempting to download file from path: {file_path}")
    if os.path.exists(file_path):
        app.logger.info(f"File found: {file_path}")
        return send_file(file_path, as_attachment=True, download_name=filename)
    else:
        app.logger.error(f"File not found: {file_path}")
        return jsonify({"error": "File not found"}), 404

if __name__ == '__main__':
    app.run(debug=True)